"""Build an editable PowerPoint of the Intake Agent end-to-end flow (16:9)."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

C = dict(
    navy=RGBColor(0x1E,0x3A,0x5F), blue=RGBColor(0x25,0x63,0xEB), blue_m=RGBColor(0xBF,0xDB,0xFE),
    blue_l=RGBColor(0xEF,0xF6,0xFF),
    green=RGBColor(0x15,0x80,0x3D), green_l=RGBColor(0xF0,0xFD,0xF4), green_m=RGBColor(0x86,0xEF,0xAC),
    amber=RGBColor(0xB4,0x53,0x09), amber_l=RGBColor(0xFF,0xFB,0xEB), amber_m=RGBColor(0xFC,0xD3,0x4D),
    purple=RGBColor(0x7C,0x3A,0xED), purple_l=RGBColor(0xF5,0xF3,0xFF), purple_m=RGBColor(0xDD,0xD6,0xFE),
    teal=RGBColor(0x0D,0x94,0x88), teal_l=RGBColor(0xF0,0xFD,0xFA), teal_m=RGBColor(0x99,0xF6,0xE4),
    gray=RGBColor(0x64,0x74,0x8B), ink=RGBColor(0x0F,0x17,0x2A), white=RGBColor(0xFF,0xFF,0xFF),
)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

def box(x, y, w, h, fill, line, lines, anchor=MSO_ANCHOR.MIDDLE, round=True, line_w=1.25, dash=None):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if round else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    tf = shp.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.06); tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.03); tf.margin_bottom = Inches(0.03)
    for i,(txt,sz,col,bold) in enumerate(lines):
        p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = txt
        r.font.size = Pt(sz); r.font.bold = bold; r.font.color.rgb = col
        r.font.name = "Segoe UI"
    return shp

def text(x, y, w, h, runs, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    p = tf.paragraphs[0]; p.alignment = align
    for txt,sz,col,bold in runs:
        r = p.add_run(); r.text = txt; r.font.size=Pt(sz); r.font.bold=bold
        r.font.color.rgb=col; r.font.name="Segoe UI"
    return tb

def arrow(x, y, w, h, fill):
    shp = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill; shp.line.fill.background()
    shp.shadow.inherit = False
    return shp

# ── Title ──
text(0.3, 0.12, 12.7, 0.5, [("Intake Agent — End-to-End Flow", 24, C['navy'], True)])
text(0.3, 0.6, 12.7, 0.35, [
    ("From any intake channel  →  an autonomous, auditable agent  →  submission to your end systems.", 12, C['gray'], False)])

# ── band labels ──
text(0.35, 1.12, 2.1, 0.3, [("INTAKE CHANNELS", 10, C['gray'], True)])
text(3.0, 1.12, 6.9, 0.3, [("THE INTAKE AGENT — AUTONOMOUS, AUDITABLE", 10, C['gray'], True)])
text(10.55, 1.12, 2.45, 0.3, [("SUBMISSION TO END SYSTEMS", 10, C['gray'], True)])

# ── INPUTS ──
inputs = [("📠  Fax","Referral packages — 3 docs"),
          ("📧  Email","Referrals & CM replies"),
          ("📞  Voice","Live AI intake line, 24/7"),
          ("🗂  Portal / Upload","Payer & partner portals")]
iy = 1.55
for t,s in inputs:
    box(0.35, iy, 2.1, 0.66, C['blue_l'], C['blue_m'],
        [(t,12,C['navy'],True),(s,8.5,C['gray'],False)])
    iy += 0.66 + 0.12

# ── big arrows ──
arrow(2.52, 2.85, 0.44, 0.55, C['blue'])
arrow(10.02, 2.85, 0.44, 0.55, C['blue'])

# ── AGENT CORE ──
box(3.0, 1.45, 6.9, 3.25, C['white'], C['blue'], [], line_w=2.0)
# tag
tagw=3.4
box(3.0+6.9/2-tagw/2, 1.34, tagw, 0.3, C['blue'], None,
    [("AGENT 1 · INTAKE — BUILT & LIVE", 9.5, C['white'], True)], round=True)

stages = [("1","Capture & Read","Reads every page or answers the call"),
          ("2","Extract","25 fields, each traced to its source"),
          ("3","Validate","Checked vs. cited CMS rules"),
          ("4","Reason & Score","Resolves conflicts; scores confidence"),
          ("5","Gaps & Outreach","Email then voice to the case manager"),
          ("6","Route-Ready","Complete record + audit trail, locked")]
sx = 3.18; sw = 1.04; sgap = 0.075; sy = 1.98; sh = 1.32
for i,(n,t,d) in enumerate(stages):
    x = sx + i*(sw+sgap)
    box(x, sy, sw, sh, C['blue_l'], C['blue_m'],
        [("STEP "+n,7,C['gray'],True),(t,10.5,C['navy'],True),(d,7.8,C['gray'],False)])
    if i < len(stages)-1:
        text(x+sw-0.02, sy, sgap+0.04, sh, [("›",13,C['blue'],True)])

# human in the loop
box(3.18, 3.46, 6.54, 0.4, C['amber_l'], C['amber_m'],
    [("👤  Human in the loop — low confidence escalates to a specialist; a person approves every outbound action", 9, C['amber'], True)])
# foundation bars
box(3.18, 3.95, 3.22, 0.5, C['purple_l'], C['purple_m'],
    [("🕸  OWL Knowledge Graph", 9.5, C['purple'], True),("24 deterministic, CMS-cited rules",7.5,C['gray'],False)])
box(6.50, 3.95, 3.22, 0.5, RGBColor(0xF1,0xF5,0xF9), RGBColor(0xCB,0xD5,0xE1),
    [("🔎  Audit trail", 9.5, C['navy'], True),("Every read, rule, decision & action logged",7.5,C['gray'],False)])

# ── END SYSTEMS ──
outs = [("📅  Scheduling","Vendor dispatch & appointments"),
        ("🏦  Billing / Claims","Claim generation & ERA posting"),
        ("🏥  Payer / EHR","Auth, eligibility, records")]
oy = 1.7
for t,s in outs:
    box(10.55, oy, 2.45, 0.82, C['green_l'], C['green_m'],
        [(t,12,C['green'],True),(s,8.5,C['gray'],False)])
    oy += 0.82 + 0.18

# ── ANALYTICS LAYER ──
box(0.35, 5.0, 12.65, 1.6, C['teal_l'], C['teal_m'], [], line_w=1.5, dash=True)
tagw=3.2
box(0.35+12.65/2-tagw/2, 4.89, tagw, 0.3, C['teal'], None,
    [("ANALYTICS & INSIGHTS LAYER", 9.5, C['white'], True)])
text(0.35, 5.18, 12.65, 0.28,
     [("Observes every episode across the whole flow — a live management view, not a separate report.", 9.5, C['teal'], False)])
metrics = [("Throughput & Latency","Volume, avg & p95 time"),
           ("Auto-Route Rate","% resolved, no human touch"),
           ("Escalation Rate","% to a specialist & why"),
           ("Confidence & KG","Avg conf, rules fired/flagged"),
           ("Top Gaps & SLA","Common missing fields, jurisdiction"),
           ("Cost per Referral","Live token + infra at volume")]
mx=0.55; mw=2.02; mgap=0.05; my=5.55; mh=0.9
for i,(t,d) in enumerate(metrics):
    box(mx+i*(mw+mgap), my, mw, mh, C['white'], C['teal_m'],
        [(t,9.5,C['teal'],True),(d,7.8,C['gray'],False)])

# ── footnote ──
text(0.3, 6.72, 12.7, 0.4,
     [("Connects via API, RPA, or file — adapts to your platform. Your data stays in your environment.   ·   Intake is Agent 1 of an 8-agent value chain, referral to payment.",
       9.5, C['gray'], False)])

desktop = os.path.join(os.path.expanduser("~"), "Desktop")
out1 = os.path.join(desktop, "Intake Agent - Flow.pptx")
out2 = os.path.join(os.path.dirname(__file__), "deliverables", "intake_flow_diagram.pptx")
prs.save(out1); prs.save(out2)
print("Saved:", out1)
print("Saved:", out2)
